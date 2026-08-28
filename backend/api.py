"""
api.py — Illomra backend (FastAPI).

This is the API layer. It reuses the engine (rag_engine.py) and exposes it
over HTTP so a Next.js/React frontend can call it.

Run it:
    uvicorn api:app --reload --port 8000
Then open http://localhost:8000/docs to try every endpoint in the browser.

Hardening (client delivery):
- Optional bearer-token auth (set APP_ACCESS_TOKEN) on everything except /health.
- Per-IP rate limit (30 req/min) + request logging middleware.
- 20 MB upload cap, SSRF guard on /link (in the engine), generic error details
  (raw exceptions go to the server log, never to the client).
"""

import json
import logging
import os
import secrets
import threading
import time
from collections import defaultdict, deque
from contextlib import asynccontextmanager

from dotenv import load_dotenv

# Load the .env sitting right next to this file BEFORE importing the engine, so
# env-driven engine config (PERSIST_DIR) sees the values no matter which folder
# uvicorn is launched from.
load_dotenv(os.path.join(os.path.dirname(__file__), ".env"))

from fastapi import (APIRouter, Depends, FastAPI, File, Header, HTTPException,
                     Request, Response, UploadFile)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel

from rag_engine import RAGEngine, extract_video_id

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s [%(name)s] %(message)s",
)
log = logging.getLogger("studymind")

MAX_UPLOAD_BYTES = 20 * 1024 * 1024  # 20 MB
UPLOAD_TOO_BIG = "File is over 20 MB — split it or upload a smaller file."

# Uploaded .docx format samples are kept here so /export can clone the ORIGINAL
# file (real header/footer/fonts/margins) instead of generating a generic one.
# On a hosted box point this at the persistent volume (e.g. /data/pattern_files).
PATTERN_DIR = os.getenv(
    "PATTERN_DIR",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "pattern_files"),
)
os.makedirs(PATTERN_DIR, exist_ok=True)

# ---- Auth modes ---------------------------------------------------------------
# google : GOOGLE_OAUTH_CLIENT_ID set — users sign in with Google; every user's
#          documents and chats are ISOLATED (owner = their Google account id).
# token  : APP_ACCESS_TOKEN set — one shared access code (single team/workspace).
# open   : neither set — local development, no auth.
GOOGLE_OAUTH_CLIENT_ID = os.getenv("GOOGLE_OAUTH_CLIENT_ID", "").strip()
APP_ACCESS_TOKEN = os.getenv("APP_ACCESS_TOKEN", "").strip()
AUTH_MODE = "google" if GOOGLE_OAUTH_CLIENT_ID else ("token" if APP_ACCESS_TOKEN else "open")

# Server-side conversation store — deployed users must not depend on one
# browser's localStorage. Point at the persistent volume in production.
CONVO_STORE = os.getenv(
    "CONVO_STORE",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "convo_store.json"),
)
_convo_lock = threading.Lock()


@asynccontextmanager
async def lifespan(app: FastAPI):
    # Warm the embedding model at boot so the FIRST question isn't slow.
    try:
        get_engine()
    except Exception:
        log.exception("Engine warm-up failed — will retry lazily on first request")
    yield


app = FastAPI(title="Illomra API", version="0.2.0", lifespan=lifespan)


# ---- Middleware --------------------------------------------------------------
# Order matters: CORS is added LAST so it wraps everything (even 429s from the
# rate limiter get CORS headers); the access log wraps the rate limiter so
# rate-limited requests still show up in the log.

RATE_LIMIT_MAX = 30       # requests ...
RATE_LIMIT_WINDOW = 60.0  # ... per this many seconds, per client IP
_rate_hits: "defaultdict[str, deque]" = defaultdict(deque)


@app.middleware("http")
async def rate_limit(request: Request, call_next):
    # /health stays open for load balancers; OPTIONS preflights are free.
    if request.url.path == "/health" or request.method == "OPTIONS":
        return await call_next(request)
    ip = request.client.host if request.client else "unknown"
    now = time.monotonic()
    hits = _rate_hits[ip]
    while hits and now - hits[0] > RATE_LIMIT_WINDOW:
        hits.popleft()
    if len(hits) >= RATE_LIMIT_MAX:
        return JSONResponse(status_code=429,
                            content={"detail": "Too many requests — slow down a little."})
    hits.append(now)
    return await call_next(request)


@app.middleware("http")
async def access_log(request: Request, call_next):
    start = time.perf_counter()
    response = await call_next(request)
    duration_ms = int((time.perf_counter() - start) * 1000)
    log.info("%s %s -> %s (%d ms)",
             request.method, request.url.path, response.status_code, duration_ms)
    return response


# CORS = permission for the frontend (running on a different origin) to call us.
app.add_middleware(
    CORSMiddleware,
    allow_origins=os.getenv("CORS_ORIGINS", "http://localhost:3000").split(","),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ---- Auth dependency -----------------------------------------------------------
# Verifies the caller and yields their OWNER id, which fences all data access:
# google mode -> the Google account's stable id ('sub'); other modes -> None
# (single-tenant: no per-user fencing, current behavior).
_google_request = None  # lazy transport (fetches+caches Google's signing certs)


def _verify_google_token(token: str) -> str:
    global _google_request
    from google.oauth2 import id_token as google_id_token
    from google.auth.transport import requests as google_requests
    if _google_request is None:
        _google_request = google_requests.Request()
    try:
        info = google_id_token.verify_oauth2_token(token, _google_request, GOOGLE_OAUTH_CLIENT_ID)
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid or expired Google sign-in — sign in again.")
    sub = str(info.get("sub", "")).strip()
    if not sub:
        raise HTTPException(status_code=401, detail="Google sign-in did not include an account id.")
    return f"g{sub}"


def current_owner(authorization: str = Header(default="")) -> "str | None":
    if AUTH_MODE == "google":
        if not authorization.startswith("Bearer "):
            raise HTTPException(status_code=401, detail="Sign in with Google to continue.")
        return _verify_google_token(authorization.removeprefix("Bearer ").strip())
    if AUTH_MODE == "token":
        expected = f"Bearer {APP_ACCESS_TOKEN}"
        if not authorization or not secrets.compare_digest(authorization, expected):
            raise HTTPException(status_code=401, detail="Missing or invalid access token.")
        return None  # shared workspace — no per-user fencing
    return None  # open/local dev


# Every route on this router requires auth. /health is registered directly on
# the app and stays open.
router = APIRouter(dependencies=[Depends(current_owner)])


# ---- Engine (one shared instance, single-process) -----------------------------
_engine = None
_engine_lock = threading.Lock()


def get_engine() -> RAGEngine:
    global _engine
    key = os.getenv("GOOGLE_API_KEY", "")
    if not key:
        raise HTTPException(
            status_code=500,
            detail="GOOGLE_API_KEY is missing. Add a line 'GOOGLE_API_KEY=your_key' to your .env file.",
        )
    if _engine is None:
        with _engine_lock:
            if _engine is None:
                _engine = RAGEngine(api_key=key)
    return _engine


# ---- Client-safe error mapping ------------------------------------------------
def _is_rate_limit_err(e) -> bool:
    s = str(e)
    return "429" in s or "RESOURCE_EXHAUSTED" in s or "quota" in s.lower() or "rate limit" in s.lower()


def _is_busy_err(e) -> bool:
    return str(e).lower().startswith("busy:")


def _friendly(e) -> str:
    if _is_rate_limit_err(e):
        return ("You've hit Gemini's free per-minute limit — wait ~30-60 seconds and try again. "
                "(Big questions use more of the limit.)")
    if _is_busy_err(e):
        return "The server is answering other questions — try again in a moment."
    return "The AI request failed — please try again in a moment."


def _ai_http_error(e) -> HTTPException:
    """Map an engine/LLM exception to a client-safe HTTPException.
    The raw exception goes to the server log, never to the client."""
    if isinstance(e, HTTPException):
        return e
    log.exception("AI request failed")
    if _is_rate_limit_err(e):
        return HTTPException(status_code=429, detail=_friendly(e))
    if _is_busy_err(e):
        return HTTPException(status_code=503, detail=_friendly(e))
    return HTTPException(status_code=502, detail=_friendly(e))


# ---- Shapes of incoming JSON (FastAPI validates these automatically) ----
class ChatIn(BaseModel):
    question: str
    history: list = []
    pattern: str = ""     # non-empty switches the engine into document-generator mode
    source: str = ""      # "" = search all documents; otherwise scope to this one
    sources: list = []    # exam-paper mode: chosen reference documents ([] = all)
    paper_opts: str = ""  # legacy free-text paper options (kept for compatibility)


class QuizIn(BaseModel):
    num_questions: int = 5
    source: str = ""
    topic: str = ""  # optional focus topic — used as the retrieval query when given
    avoid: list = []  # stems of recently asked questions, to prevent repeats


class FlashcardsIn(BaseModel):
    num_cards: int = 8
    source: str = ""


class ConvosIn(BaseModel):
    convos: list = []
    current: str = ""


class SummaryIn(BaseModel):
    source: str = ""


class DeleteDocIn(BaseModel):
    source: str


class UrlIn(BaseModel):
    url: str


class ExportIn(BaseModel):
    text: str
    format: str = "pdf"
    pattern_id: str = ""  # id of a stored .docx sample — export clones that file
    layout: dict = {}     # visual spec from a photo sample: font / centered header / footer


# A tiny adapter so the engine's process_file (built for Streamlit uploads)
# also accepts FastAPI uploads. The engine only needs .name and .read().
class _UploadShim:
    def __init__(self, filename: str, data: bytes):
        self.name = filename
        self._data = data

    def read(self) -> bytes:
        return self._data


def _parse_inline_md(text: str):
    """Splits text into tuples of (substring, is_bold, is_italic, is_code)."""
    import re
    tokens = []
    pattern = re.compile(r'(\*\*.*?\*\*|\*.*?\*|`.*?`)')
    parts = pattern.split(text)
    for p in parts:
        if not p:
            continue
        if p.startswith('**') and p.endswith('**'):
            tokens.append((p[2:-2], True, False, False))
        elif p.startswith('*') and p.endswith('*'):
            tokens.append((p[1:-1], False, True, False))
        elif p.startswith('`') and p.endswith('`'):
            tokens.append((p[1:-1], False, False, True))
        else:
            tokens.append((p, False, False, False))
    return tokens


def _add_docx_p(doc, text: str, style=None, align=None, heading_level=None):
    from docx.shared import Pt, RGBColor
    if heading_level:
        p = doc.add_heading('', level=heading_level)
    elif style:
        p = doc.add_paragraph(style=style)
    else:
        p = doc.add_paragraph()

    if align:
        p.alignment = align

    tokens = _parse_inline_md(text)
    for content, bold, italic, code in tokens:
        run = p.add_run(content)
        if bold:
            run.bold = True
        if italic:
            run.italic = True
        if code:
            run.font.name = 'Consolas'
            run.font.size = Pt(9.5)
            run.font.color.rgb = RGBColor(180, 40, 40)
    return p


def _fill_docx(doc, text: str, layout: dict, center_rest: bool = False):
    """Write formatted markdown paper text into a python-docx Document."""
    import re
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches
    center_n = int((layout or {}).get("center_top_lines", 0) or 0)
    lines = text.split("\n")
    line_no = 0
    i = 0

    while i < len(lines):
        raw = lines[i]
        st = raw.strip()
        if not st:
            doc.add_paragraph("")
            i += 1
            continue

        line_no += 1
        align = WD_ALIGN_PARAGRAPH.CENTER if (line_no <= center_n or center_rest) else None

        # Check for Markdown Table
        if st.startswith('|') and st.endswith('|'):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                table_lines.append(lines[i].strip())
                i += 1
            rows = []
            for t_line in table_lines:
                cells = [c.strip() for c in t_line.strip('|').split('|')]
                if all(re.match(r'^:?-+:?$', c) for c in cells if c):
                    continue
                rows.append(cells)
            if rows:
                col_count = max(len(r) for r in rows)
                tbl = doc.add_table(rows=len(rows), cols=col_count)
                tbl.style = 'Table Grid'
                for r_idx, r_data in enumerate(rows):
                    for c_idx, cell_text in enumerate(r_data):
                        if c_idx < col_count:
                            cell_p = tbl.cell(r_idx, c_idx).paragraphs[0]
                            for content, bold, italic, code in _parse_inline_md(cell_text):
                                run = cell_p.add_run(content)
                                if r_idx == 0 or bold:
                                    run.bold = True
            continue

        # Headings & Lists
        if st.startswith("### "):
            _add_docx_p(doc, st[4:], heading_level=3, align=align)
        elif st.startswith("## "):
            _add_docx_p(doc, st[3:], heading_level=2, align=align)
        elif st.startswith("# "):
            _add_docx_p(doc, st[2:], heading_level=1, align=align)
        elif st.startswith(("- ", "* ", "+ ")):
            _add_docx_p(doc, st[2:], style="List Bullet", align=align)
        elif re.match(r'^\d+\.\s+', st):
            clean_text = re.sub(r'^\d+\.\s+', '', st)
            _add_docx_p(doc, clean_text, style="List Number", align=align)
        elif st.startswith("> "):
            p = _add_docx_p(doc, st[2:], align=align)
            p.paragraph_format.left_indent = Inches(0.4)
        else:
            _add_docx_p(doc, st, align=align)
        i += 1


def _paper_to_bytes(text: str, fmt: str, layout: dict = None, pattern_path: str = None):
    """Convert a markdown paper into pdf / docx / pptx bytes."""
    import re
    layout = layout or {}
    if fmt == "pdf":
        text = text.translate(str.maketrans({
            "—": "-", "–": "-", "‘": "'", "’": "'",
            "“": '"', "”": '"', "…": "...", " ": " ",
            "•": "-", "→": "->", "×": "x", "−": "-",
        }))
        try:
            text.encode("latin-1")
        except UnicodeEncodeError:
            raise ValueError("PDF export can't render Urdu or special characters yet — download as Word instead (full support).")
        from fpdf import FPDF
        from fpdf.enums import XPos, YPos

        family = "Times" if layout.get("font") == "serif" else "Helvetica"
        footer_text = str(layout.get("footer", "") or "")

        class PaperPDF(FPDF):
            def footer(self):
                self.set_y(-13)
                self.set_font(family, "I", 9)
                left = footer_text[:80]
                self.cell(self.epw / 2, 6, left, align="L")
                self.cell(self.epw / 2, 6, f"Page {self.page_no()}", align="R")

        pdf = PaperPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=18)
        center_n = int(layout.get("center_top_lines", 0) or 0)
        lines = text.split("\n")
        line_no = 0
        i = 0

        def _line(t, size, style="", align="L"):
            pdf.set_font(family, style, size)
            pdf.set_x(pdf.l_margin)
            # Strip markdown formatting symbols for clean PDF rendering
            clean_t = t.replace("**", "").replace("*", "").replace("`", "")
            pdf.multi_cell(pdf.epw, size * 0.55, clean_t, align=align,
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        while i < len(lines):
            raw = lines[i]
            st = raw.strip()
            if not st:
                pdf.ln(3)
                i += 1
                continue
            line_no += 1
            centered = "C" if line_no <= center_n else "L"

            # Table handling in PDF
            if st.startswith('|') and st.endswith('|'):
                table_lines = []
                while i < len(lines) and lines[i].strip().startswith('|'):
                    table_lines.append(lines[i].strip())
                    i += 1
                rows = []
                for t_line in table_lines:
                    cells = [c.strip().replace("**", "") for c in t_line.strip('|').split('|')]
                    if all(re.match(r'^:?-+:?$', c) for c in cells if c):
                        continue
                    rows.append(cells)
                if rows:
                    col_count = max(len(r) for r in rows)
                    col_width = pdf.epw / col_count
                    pdf.set_font(family, "B", 10)
                    for r_idx, r_data in enumerate(rows):
                        style = "B" if r_idx == 0 else ""
                        pdf.set_font(family, style, 9.5)
                        for c_idx in range(col_count):
                            cell_val = r_data[c_idx] if c_idx < len(r_data) else ""
                            pdf.cell(col_width, 6, cell_val[:30], border=1)
                        pdf.ln()
                continue

            if st.startswith("### "):
                _line(st[4:], 12, "B", centered)
            elif st.startswith("## "):
                _line(st[3:], 14, "B", centered)
            elif st.startswith("# "):
                _line(st[2:], 16, "B", centered)
            elif st.startswith(("- ", "* ", "+ ")):
                _line("   • " + st[2:], 11, "", centered)
            elif re.match(r'^\d+\.\s+', st):
                _line("   " + st, 11, "", centered)
            else:
                _line(st, 11, "B" if line_no <= center_n else "", centered)
            i += 1
        return bytes(pdf.output()), "application/pdf", "pdf"
    if fmt == "docx":
        import io
        from docx import Document as Docx
        if pattern_path and os.path.exists(pattern_path):
            # EXACT-format path: clone the teacher's own sample file. Its header,
            # footer, margins, styles, and fonts stay untouched — we only replace
            # the body content with the new paper.
            doc = Docx(pattern_path)
            for tbl in list(doc.tables):
                tbl._element.getparent().remove(tbl._element)
            for p in list(doc.paragraphs):
                p._element.getparent().remove(p._element)
            _fill_docx(doc, text, layout)
        else:
            doc = Docx()
            if layout.get("font") == "serif":
                doc.styles["Normal"].font.name = "Times New Roman"
            _fill_docx(doc, text, layout)
            footer_text = str(layout.get("footer", "") or "")
            if footer_text:
                try:
                    doc.sections[0].footer.paragraphs[0].text = footer_text[:200]
                except Exception:
                    log.exception("Could not write docx footer")
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"
    if fmt == "pptx":
        import io
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()

        # Split text into slide groups based on headings
        slides_data = []
        current_slide = {"title": "Presentation", "content": []}

        for raw_line in text.split("\n"):
            st = raw_line.strip()
            if not st:
                continue
            if st.startswith("# ") or st.startswith("## "):
                if current_slide["content"] or current_slide["title"] != "Presentation":
                    slides_data.append(current_slide)
                title_clean = st.lstrip("#").replace("**", "").strip()
                current_slide = {"title": title_clean, "content": []}
            else:
                clean = st.replace("**", "").replace("*", "").strip()
                current_slide["content"].append(clean)

        if current_slide["content"] or current_slide["title"]:
            slides_data.append(current_slide)

        if not slides_data:
            slides_data = [{"title": "Presentation", "content": [text]}]

        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

        # Slide 1: Title Slide
        first_slide = prs.slides.add_slide(title_layout)
        first_slide.shapes.title.text = slides_data[0]["title"]
        if first_slide.placeholders[1]:
            first_slide.placeholders[1].text = "Generated by StudyMind AI"

        # Subsequent slides
        for slide_item in slides_data[1:]:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_item["title"]
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            for idx, bullet_text in enumerate(slide_item["content"][:8]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = bullet_text
                p.font.size = Pt(14)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"
    raise ValueError("Unsupported format")


# ---- Routes -------------------------------------------------------------------
@app.get("/health")
def health():
    try:
        eng = get_engine()
    except HTTPException:
        return {"status": "degraded", "documents": 0, "chunks": 0, "store_ok": False}
    return {
        "status": "ok",
        "auth": AUTH_MODE,
        "documents": eng.doc_count,
        "chunks": eng.chunk_count,
        "store_ok": not eng.store_error,
    }


@router.post("/upload")
async def upload(file: UploadFile = File(...), owner: "str | None" = Depends(current_owner)):
    data = await file.read()
    if len(data) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail=UPLOAD_TOO_BIG)
    try:
        chunks = get_engine().process_file(_UploadShim(file.filename, data), owner=owner)
    except HTTPException:
        raise
    except Exception as e:
        # Photo transcription goes through Gemini — surface limit/busy errors honestly.
        if _is_rate_limit_err(e) or _is_busy_err(e):
            raise _ai_http_error(e)
        log.exception("Upload failed for %r", file.filename)
        raise HTTPException(status_code=400,
                            detail="Couldn't read that file — supported: PDF, Word, PPT, TXT, MD, CSV, or a photo (JPG/PNG).")
    return {"filename": file.filename, "chunks": chunks, "stats": get_engine().get_stats(owner)}


@router.post("/link")
def link(body: UrlIn, owner: "str | None" = Depends(current_owner)):
    """Smart link handler: YouTube link -> transcript, anything else -> web page."""
    url = body.url.strip()
    try:
        if extract_video_id(url):
            chunks = get_engine().process_youtube_url(url, owner=owner)
        else:
            chunks = get_engine().process_url(url, owner=owner)
    except HTTPException:
        raise
    except ValueError as e:
        # Our own guard messages (bad YouTube link, SSRF rejection) — safe to show.
        raise HTTPException(status_code=400, detail=str(e))
    except Exception:
        log.exception("Link ingestion failed for %r", url)
        raise HTTPException(status_code=400,
                            detail="Couldn't read that link — the site may block bots or need a login.")
    return {"chunks": chunks, "stats": get_engine().get_stats(owner)}


@router.post("/chat-stream")
def chat_stream(body: ChatIn, owner: "str | None" = Depends(current_owner)):
    if not body.question.strip():
        raise HTTPException(status_code=400, detail="Ask a question first.")

    def gen():
        try:
            for part in get_engine().query_stream(
                body.question,
                history=body.history,
                pattern=body.pattern,
                source=body.source,
                paper_opts=body.paper_opts,
                ref_sources=[s for s in body.sources if isinstance(s, str) and s] or None,
                owner=owner,
            ):
                yield json.dumps(part) + "\n"
        except Exception as e:
            log.exception("chat-stream failed")
            yield json.dumps({"done": True, "error": _friendly(e),
                              "sources": [], "confidence": 0, "usage": {}}) + "\n"

    return StreamingResponse(gen(), media_type="application/x-ndjson", headers={
        "X-Accel-Buffering": "no",
        "Cache-Control": "no-cache",
        "Connection": "keep-alive",
    })


@router.post("/quiz")
def quiz(body: QuizIn, owner: "str | None" = Depends(current_owner)):
    try:
        return {"questions": get_engine().generate_quiz(
            num_questions=body.num_questions, source=body.source, topic=body.topic,
            avoid=[a for a in body.avoid if isinstance(a, str)][:20], owner=owner)}
    except Exception as e:
        raise _ai_http_error(e)


@router.post("/flashcards")
def flashcards(body: FlashcardsIn, owner: "str | None" = Depends(current_owner)):
    try:
        return {"cards": get_engine().generate_flashcards(
            num_cards=min(max(body.num_cards, 3), 20), source=body.source, owner=owner)}
    except Exception as e:
        raise _ai_http_error(e)


@router.post("/summary")
def summary(body: SummaryIn = SummaryIn(), owner: "str | None" = Depends(current_owner)):
    try:
        return get_engine().summarize(source=body.source, owner=owner)
    except Exception as e:
        raise _ai_http_error(e)


@router.get("/documents")
def documents(owner: "str | None" = Depends(current_owner)):
    """List the caller's indexed documents (powers the doc manager)."""
    return {"documents": get_engine().list_documents(owner)}


@router.post("/documents/delete")
def delete_document(body: DeleteDocIn, owner: "str | None" = Depends(current_owner)):
    try:
        stats = get_engine().delete_document(body.source, owner=owner)
    except HTTPException:
        raise
    except Exception:
        log.exception("Delete failed for %r", body.source)
        raise HTTPException(status_code=400, detail="Could not delete that document — try again.")
    return {"status": "deleted", "stats": stats}


@router.post("/extract")
async def extract(file: UploadFile = File(...)):
    """Extract text from an uploaded sample/pattern paper (not indexed as content).
    - Photos: Gemini vision transcribes AND reads the visual layout (font,
      centered header block, footer) so exports can mirror the look.
    - .docx samples: the raw file is also stored, so /export can clone it for an
      EXACT-format paper (real header/footer/fonts preserved)."""
    import hashlib
    import tempfile
    data = await file.read()
    if len(data) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail=UPLOAD_TOO_BIG)
    ext = "." + (file.filename or "x.txt").split(".")[-1].lower()
    if ext in RAGEngine.IMAGE_EXTS:
        try:
            text, layout = get_engine().image_to_text(data, with_layout=True)
        except Exception as e:
            raise _ai_http_error(e)
        return {"text": text[:15000], "layout": layout}
    pattern_id = ""
    if ext == ".docx":
        pattern_id = hashlib.sha256(data).hexdigest()[:16]
        try:
            with open(os.path.join(PATTERN_DIR, f"{pattern_id}.docx"), "wb") as f:
                f.write(data)
        except Exception:
            log.exception("Could not store pattern file")
            pattern_id = ""
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(data)
        path = tmp.name
    try:
        docs = RAGEngine._load_any(path, ext)
        text = "\n\n".join(d.page_content for d in docs)
    except Exception:
        log.exception("Extract failed for %r", file.filename)
        raise HTTPException(status_code=400,
                            detail="Couldn't read that file — supported: PDF, Word, PPT, TXT, MD, CSV, or a photo.")
    finally:
        os.unlink(path)
    return {"text": text[:15000], "pattern_id": pattern_id}


@router.post("/export")
def export(body: ExportIn):
    import re as _re
    pattern_path = None
    if body.pattern_id and _re.fullmatch(r"[0-9a-f]{16}", body.pattern_id):
        candidate = os.path.join(PATTERN_DIR, f"{body.pattern_id}.docx")
        if os.path.exists(candidate):
            pattern_path = candidate
    try:
        data, mime, ext = _paper_to_bytes(body.text, body.format,
                                          layout=body.layout, pattern_path=pattern_path)
    except ValueError as e:
        # Our own messages (Urdu/PDF limitation, unsupported format) — safe to show.
        raise HTTPException(status_code=400, detail=str(e))
    except Exception:
        log.exception("Export failed")
        raise HTTPException(status_code=400, detail="Export failed — try a different format.")
    return Response(content=data, media_type=mime,
                    headers={"Content-Disposition": f'attachment; filename="studymind-paper.{ext}"'})


@router.get("/stats")
def stats(owner: "str | None" = Depends(current_owner)):
    return get_engine().get_stats(owner)


def _convo_path(owner: "str | None") -> str:
    """One conversation file per owner. Single-tenant modes keep the legacy path
    so existing chats survive; Google users each get their own file."""
    if not owner:
        return CONVO_STORE
    safe = "".join(ch for ch in owner if ch.isalnum())[:64] or "unknown"
    d = CONVO_STORE + ".d"
    os.makedirs(d, exist_ok=True)
    return os.path.join(d, f"{safe}.json")


@router.get("/convos")
def get_convos(owner: "str | None" = Depends(current_owner)):
    """The saved conversation list — chats follow the account, not one browser."""
    try:
        with open(_convo_path(owner), "r", encoding="utf-8") as f:
            data = json.load(f)
        return {"convos": data.get("convos", []), "current": data.get("current", "")}
    except Exception:
        return {"convos": [], "current": ""}


@router.put("/convos")
def put_convos(body: ConvosIn, owner: "str | None" = Depends(current_owner)):
    data = {"convos": body.convos[:300], "current": body.current}
    raw = json.dumps(data, ensure_ascii=False)
    if len(raw) > 8_000_000:
        raise HTTPException(status_code=413,
                            detail="Chat history is too large to save — delete some old chats.")
    path = _convo_path(owner)
    with _convo_lock:
        tmp = path + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            f.write(raw)
        os.replace(tmp, path)
    return {"status": "saved"}


@router.get("/usage")
def usage():
    """Per-model free-tier usage for the UI's limit bar (our own counts —
    Google exposes no remaining-quota API, so daily_limit values are estimates).
    The Gemini key is shared, so this is workspace-wide by design."""
    return get_engine().usage_snapshot()


@router.post("/reset")
def reset(owner: "str | None" = Depends(current_owner)):
    """Clear materials — only the caller's own in Google mode; everything otherwise."""
    get_engine().reset(owner=owner)
    return {"status": "reset"}


app.include_router(router)
